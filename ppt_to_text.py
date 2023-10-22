# convert a ppt file into plain text
from pptx import Presentation

import gpt

import re # for cleaning text

import os

def extract_text_from_powerpoint(path):
    # Check if the file exists
    if not os.path.exists(path):
        raise FileNotFoundError(f"The file '{path}' does not exist.")

    # Determine the file type based on the extension
    file_extension = os.path.splitext(path)[-1].lower()

    if file_extension == ".pptx":
        # Handle PowerPoint files (PPTX)
        text = ""
        presentation = Presentation(path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        raise ValueError("Unsupported file format. Only .pptx and .pdf files are supported.")



def clean_text(sample_text):
    # Remove line numbers (e.g., "9 of 59")
    sample_text = re.sub(r'\d+\s+of\s+\d+', '', sample_text)
    
    # Remove extra whitespace and newlines
    sample_text = ' '.join(sample_text.split())
    
    # Remove leading and trailing spaces
    sample_text = sample_text.strip()
    
    return sample_text

if __name__ == '__main__':
    import os

    path = os.path.join(os.getcwd(), 'testppt.pptx')

    plain_text = extract_text_from_powerpoint(path)
    
    # print(plain_text)

    # print(clean_text(plain_text))

    b = gpt.gpt(gpt.keys.gpt_api_key)
    print(b.chat("make a list of flash cards for this " + clean_text(plain_text)))

    

    